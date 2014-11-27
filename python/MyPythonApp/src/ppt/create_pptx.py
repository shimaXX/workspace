# coding: utf-8

from pptx import Presentation
from pptx.util import Inches, Px, Pt

def create_pptx(region_list, caluculateion_type='live'):
    img_basepath = 'C:/workspace/py/MyPythonApp/src/viewingScore/create_graph/%s/' % caluculateion_type
    SLD_LAYOUT_TITLE_AND_CONTENT = 1
    
    gender_list = ('男性', '女性')
    week_list = ('','weekday_','weekend_')
    time_list=('5-9','9-19','19-23','23-29')
    img_gender_name_list = ('r_m_', 'r_fm_')
    img_name_tips=('20to34', '35to49', 'over50')
    #'output_'+region+
    prs = Presentation()
    for region in region_list:
        # create slide to describe region
        blank_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(blank_slide_layout)
        shapes = slide.shapes
        
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
    
        title_shape.text = region # title
    
        # create slide with graph to describe hole users having properties at that region 
        blank_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(blank_slide_layout)
        shapes = slide.shapes
        
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        title_shape.text = region+'_個人全体' # title
    
        # paste image
        width = height = Px(600)
    
        base_left = 1.8
        base_top = 1.25
    
        img_path = unicode(img_basepath+'output_'+region+'_r_individuals.png')
        print img_path
        left = Inches(base_left)
        top = Inches(base_top)
        pic = slide.shapes.add_picture(img_path, left, top, width, height)
    
        for _week in week_list:
            if _week=='':
                for i, gender in enumerate(gender_list):
                    blank_slide_layout = prs.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
                    slide = prs.slides.add_slide(blank_slide_layout)
                    shapes = slide.shapes
                
                    title_shape = shapes.title
                    body_shape = shapes.placeholders[1]
                
                    title_shape.text = region+'_'+gender # title
                
                    # paste image
                    img_width = img_height = Px(300)
                    # image base
                    base_left = 0.5
                    base_top = 2
                    # text base
                    txt_base_left = 1.2
                    txt_base_top = 1.4
                    for img_idx in xrange(len(img_name_tips)):
                        left_inch = base_left+img_idx*3
                        top_inch = base_top+img_idx*1
                
                        img_path = unicode(img_basepath+'output_'+region+'_'+_week+img_gender_name_list[i]+img_name_tips[img_idx]+'.png')
                        left = Inches(left_inch)
                        top = Inches(top_inch)
                        print img_path
                        pic = slide.shapes.add_picture(img_path, left, top, img_width, img_height)
                
                        # write text
                        left_inch = txt_base_left+img_idx*3
                        top_inch = txt_base_top+img_idx*1
                
                        left = Inches(left_inch)
                        top = Inches(top_inch)
                        width = height = Inches(1)
                        txBox = slide.shapes.add_textbox(left, top, width, height)
                        tf = txBox.textframe
                        p = tf.add_paragraph()
                        p.text = gender+img_name_tips[img_idx]
                        p.font.size = Pt(24)
                        p.font.bold = True

            for _time in time_list:
                if _week!='':
                    for i, gender in enumerate(gender_list):
                        blank_slide_layout = prs.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
                        slide = prs.slides.add_slide(blank_slide_layout)
                        shapes = slide.shapes
                    
                        title_shape = shapes.title
                        body_shape = shapes.placeholders[1]
                    
                        title_shape.text = region+'_'+_week+_time+'_'+gender # title
                    
                        # paste image
                        img_width = img_height = Px(300)
                        # image base
                        base_left = 0.5
                        base_top = 2
                        # text base
                        txt_base_left = 1.2
                        txt_base_top = 1.4
                        for img_idx in xrange(len(img_name_tips)):
                            left_inch = base_left+img_idx*3
                            top_inch = base_top+img_idx*1
                    
                            img_path = unicode(img_basepath+'output_'+region+'_'+_week+img_gender_name_list[i]+img_name_tips[img_idx]+'_'+_time+'.png')
                            left = Inches(left_inch)
                            top = Inches(top_inch)
                            print img_path
                            pic = slide.shapes.add_picture(img_path, left, top, img_width, img_height)
                    
                            # write text
                            left_inch = txt_base_left+img_idx*3
                            top_inch = txt_base_top+img_idx*1
                    
                            left = Inches(left_inch)
                            top = Inches(top_inch)
                            width = height = Inches(1)
                            txBox = slide.shapes.add_textbox(left, top, width, height)
                            tf = txBox.textframe
                            p = tf.add_paragraph()
                            p.text = gender+img_name_tips[img_idx]
                            p.font.size = Pt(24)
                            p.font.bold = True                
    prs.save('%s.pptx' % caluculateion_type) 
    
if __name__=='__main__':
    region_list = ('関東','関西','名古屋')
    try:
        create_pptx(region_list, caluculateion_type='live')
    except Exception, e:
        print e
# coding: utf-8

from pptx import Presentation
from pptx.util import Inches, Px, Pt

img_basepath = 'my img directory'
img_path_list = (('age1.jpg','age2.jpg','age3.jpg'),('age1.jpg','age2.jpg','age3.jpg')) # idxで取得出来る用にlist化
classification_list = ('male','female')
age_range_list = ('2-34','35-49','50-')
SLD_LAYOUT_TITLE_AND_CONTENT = 1

prs = Presentation()
for i, class_name in enumerate(classification_list):
    blank_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(blank_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[SLD_LAYOUT_TITLE_AND_CONTENT]
    
    title_shape.text = class_name # title
    
    # paste image
    width = height = Px(300)

    base_left = 0.5
    base_top = 2
    for img_idx in xrange(len(img_path_list[i])):
        left_inch = base_left+img_idx*3
        top_inch = base_top+img_idx*1
        
        img_path = class_name+'_'+img_path_list[i][img_idx]
        left = Inches(left_inch)
        top = Inches(top_inch)
        pic = slide.shapes.add_picture(img_path, left, top, width, height)
        
    # write text
    base_left = 1.2
    base_top = 1.4
    for age_idx in xrange(len(img_path_list[i])):
        left_inch = base_left+age_idx*3
        top_inch = base_top+age_idx*1
        
        left = Inches(left_inch)
        top = Inches(top_inch)
        width = height = Inches(1)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.textframe
        p = tf.add_paragraph()
        p.text = class_name+age_range_list[age_idx]
        p.font.size = Pt(24)
        p.font.bold = True

prs.save('test.pptx')